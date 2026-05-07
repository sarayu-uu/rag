"""
File purpose:
- Contains all ingestion loaders for supported input types.
- Each loader converts its input source into plain text for downstream RAG processing.
"""

from pathlib import Path
import subprocess
import tempfile
from typing import Any, Union

# PDF fallback
import fitz  # PyMuPDF

# DOCX fallback
from docx import Document

# PPTX fallback
from pptx import Presentation

# Image OCR fallback
from PIL import Image
import pytesseract

# CSV fallback
import pandas as pd

# Web fallback
import requests
from bs4 import BeautifulSoup

# JSON/XML fallback
import json
import xml.etree.ElementTree as ET


def _extract_page_number(metadata: dict[str, Any]) -> int | None:
    for key in ("page", "page_number"):
        value = metadata.get(key)
        if isinstance(value, int):
            return value + 1 if key == "page" and value >= 0 else value
        if isinstance(value, str) and value.isdigit():
            num = int(value)
            return num + 1 if key == "page" and num >= 0 else num
    return None


def _documents_to_text(documents: list[Any]) -> str:
    parts: list[str] = []
    for doc in documents:
        value = str(getattr(doc, "page_content", "") or "").strip()
        if value:
            parts.append(value)
    return "\n".join(parts)


def _load_documents_with_langchain(loader_name: str, *args: Any, **kwargs: Any) -> list[Any] | None:
    try:
        from langchain_community import document_loaders as lc_loaders
    except Exception:
        return None

    loader_cls = getattr(lc_loaders, loader_name, None)
    if loader_cls is None:
        return None

    try:
        loader = loader_cls(*args, **kwargs)
        if hasattr(loader, "load"):
            return loader.load()
    except Exception:
        return None
    return None


# Runs powershell office conversion.
def _run_powershell_office_conversion(
    script: str,
    source_path: Path,
    output_path: Path,
) -> None:
    try:
        subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-NonInteractive",
                "-Command",
                script,
                str(source_path),
                str(output_path),
            ],
            capture_output=True,
            text=True,
            check=True,
            timeout=120,
        )
    except FileNotFoundError as exc:
        raise ValueError("PowerShell is not available for legacy Office file conversion.") from exc
    except subprocess.TimeoutExpired as exc:
        raise ValueError("Legacy Office file conversion timed out.") from exc
    except subprocess.CalledProcessError as exc:
        details = exc.stderr.strip() or exc.stdout.strip() or "Unknown conversion error."
        raise ValueError(f"Legacy Office file conversion failed: {details}") from exc


# Loads pdf.
def load_pdf(file_path: Union[str, Path]) -> str:
    sections = load_pdf_sections(file_path)
    return "\n".join(section["text"] for section in sections)


# Loads pdf sections.
def load_pdf_sections(file_path: Union[str, Path]) -> list[dict[str, Union[int, str]]]:
    path = str(file_path)

    docs = _load_documents_with_langchain("PyPDFLoader", path)
    if docs is not None:
        sections: list[dict[str, Union[int, str]]] = []
        for idx, doc in enumerate(docs, start=1):
            text = str(getattr(doc, "page_content", "") or "")
            metadata = getattr(doc, "metadata", {}) or {}
            page_number = _extract_page_number(metadata) or idx
            sections.append({"page_number": page_number, "text": text})
        return sections

    doc = fitz.open(file_path)
    sections = []
    for page_index, page in enumerate(doc, start=1):
        sections.append({"page_number": page_index, "text": page.get_text()})
    doc.close()
    return sections


# Loads doc.
def load_doc(file_path: Union[str, Path]) -> str:
    source_path = Path(file_path)
    with tempfile.TemporaryDirectory() as temp_dir:
        converted_path = Path(temp_dir) / f"{source_path.stem}.docx"
        script = (
            "$ErrorActionPreference='Stop';"
            "$sourcePath=$args[0];"
            "$outputPath=$args[1];"
            "$word=$null;"
            "$document=$null;"
            "try {"
            "  $word=New-Object -ComObject Word.Application;"
            "  $word.Visible=$false;"
            "  $document=$word.Documents.Open($sourcePath);"
            "  $document.SaveAs([ref]$outputPath,[ref]16);"
            "} finally {"
            "  if ($document -ne $null) { $document.Close() }"
            "  if ($word -ne $null) { $word.Quit() }"
            "}"
        )
        _run_powershell_office_conversion(script, source_path, converted_path)
        return load_docx(converted_path)


# Loads docx.
def load_docx(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain("Docx2txtLoader", path)
    if docs is not None:
        return _documents_to_text(docs)

    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


# Loads ppt.
def load_ppt(file_path: Union[str, Path]) -> str:
    source_path = Path(file_path)
    with tempfile.TemporaryDirectory() as temp_dir:
        converted_path = Path(temp_dir) / f"{source_path.stem}.pptx"
        script = (
            "$ErrorActionPreference='Stop';"
            "$sourcePath=$args[0];"
            "$outputPath=$args[1];"
            "$powerpoint=$null;"
            "$presentation=$null;"
            "try {"
            "  $powerpoint=New-Object -ComObject PowerPoint.Application;"
            "  $presentation=$powerpoint.Presentations.Open($sourcePath,$false,$false,$false);"
            "  $presentation.SaveAs($outputPath,24);"
            "} finally {"
            "  if ($presentation -ne $null) { $presentation.Close() }"
            "  if ($powerpoint -ne $null) { $powerpoint.Quit() }"
            "}"
        )
        _run_powershell_office_conversion(script, source_path, converted_path)
        return load_pptx(converted_path)


# Loads pptx.
def load_pptx(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain("UnstructuredPowerPointLoader", path)
    if docs is not None:
        return _documents_to_text(docs)

    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Loads image.
def load_image(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain("UnstructuredImageLoader", path)
    if docs is not None:
        return _documents_to_text(docs)

    img = Image.open(file_path)
    text = pytesseract.image_to_string(img)
    return text


# Loads csv.
def load_csv(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain("CSVLoader", path)
    if docs is not None:
        return _documents_to_text(docs)

    df = pd.read_csv(file_path)
    text = ""
    for _, row in df.iterrows():
        row_text = ", ".join([f"{col}: {row[col]}" for col in df.columns])
        text += row_text + "\n"
    return text


# Loads txt.
def load_txt(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain("TextLoader", path, encoding="utf-8")
    if docs is not None:
        return _documents_to_text(docs)

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# Loads web.
def load_web(url: str) -> str:
    docs = _load_documents_with_langchain("WebBaseLoader", web_paths=[url])
    if docs is not None:
        return _documents_to_text(docs)

    response = requests.get(url, timeout=20)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")

    for script in soup(["script", "style"]):
        script.decompose()

    return soup.get_text(separator="\n")


# Loads json.
def load_json(file_path: Union[str, Path]) -> str:
    path = str(file_path)
    docs = _load_documents_with_langchain(
        "JSONLoader",
        path,
        jq_schema=".",
        text_content=False,
    )
    if docs is not None:
        return _documents_to_text(docs)

    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2)


# Loads xml.
def load_xml(file_path: Union[str, Path]) -> str:
    tree = ET.parse(file_path)
    root = tree.getroot()

    parts: list[str] = []

    # Walks through XML elements and collects their text and attributes.
    def walk(element: ET.Element, path: str) -> None:
        current_path = f"{path}/{element.tag}" if path else element.tag

        text = (element.text or "").strip()
        if text:
            parts.append(f"{current_path}: {text}")

        for key, value in element.attrib.items():
            parts.append(f"{current_path}[@{key}]: {value}")

        for child in list(element):
            walk(child, current_path)

    walk(root, "")
    return "\n".join(parts)
